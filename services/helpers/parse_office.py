"""Microsoft Office parsers — Word (.docx/.doc) and PowerPoint (.pptx).

Packaged together because they share the python-docx / python-pptx
dependency family. Registers text, image, and (for Word) tabular extraction
for each. All heavy imports are lazy so the module degrades to "not installed"
rather than failing the parser-discovery scan.
"""


dependencies_files = []
dependencies_pip = ['Pillow', 'pandas', 'python-docx', 'python-pptx']

import logging
import time
from pathlib import Path

from plugins.services.helpers.ParseResult import ParseResult
from plugins.services.helpers import parser_registry as registry
from plugins.services.helpers.parsing_utils import clean_text, max_chars

logger = logging.getLogger("ParseOffice")

# Safety cap for pathological DOCX tables (matches the tabular package).
DEFAULT_MAX_TABLE_ROWS = 100_000


# ===================================================================
# DOCX
# ===================================================================

def parse_docx_text(path: str, config: dict, services: dict = None) -> ParseResult:
    """Extract text from a Word document. Detects embedded images."""
    try:
        from docx import Document
    except ImportError:
        logger.debug("python-docx not installed")
        return ParseResult.failed("python-docx not installed", modality="text")

    try:
        t0 = time.time()
        limit = max_chars(config)
        doc = Document(path)

        paragraphs = []
        current_len = 0
        for para in doc.paragraphs:
            paragraphs.append(para.text)
            current_len += len(para.text)
            if current_len > limit:
                break

        content = clean_text("\n".join(paragraphs)[:limit])

        also_contains = []
        image_count = 0
        has_tables = len(doc.tables) > 0
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                image_count += 1

        metadata = {
            "char_count": len(content),
            "paragraph_count": len(doc.paragraphs),
            "image_count": image_count,
            "has_images": image_count > 0,
            "has_tables": has_tables,
            "table_count": len(doc.tables),
        }

        if image_count > 0:
            also_contains.append("image")
        if has_tables:
            also_contains.append("tabular")

        logger.debug(
            f"DOCX parsed: {Path(path).name} — {len(doc.paragraphs)} paragraphs, "
            f"{len(content)} chars in {time.time() - t0:.2f}s"
        )
        return ParseResult(
            modality="text",
            output=content,
            metadata=metadata,
            also_contains=also_contains,
        )
    except Exception as e:
        logger.debug(f"Failed to parse {path}: {e}")
        return ParseResult.failed(str(e), modality="text")


registry.register([".docx", ".doc"], "text", parse_docx_text)


def parse_docx_image(path: str, config: dict, services: dict = None) -> ParseResult:
    """Extract embedded images from a DOCX as PIL.Image objects."""
    try:
        from docx import Document
        from PIL import Image
        import io
    except ImportError as e:
        logger.debug(f"Missing dependency: {e}")
        return ParseResult.failed(f"Missing dependency: {e}", modality="image")

    try:
        doc = Document(path)
        images = []
        max_images = config.get("max_images", 50)

        for rel in doc.part.rels.values():
            if len(images) >= max_images:
                break
            if "image" in rel.reltype:
                try:
                    image_data = rel.target_part.blob
                    img = Image.open(io.BytesIO(image_data))
                    images.append(img)
                except Exception:
                    continue

        if not images:
            return ParseResult.failed(
                "No extractable images found in DOCX", modality="image"
            )

        return ParseResult(
            modality="image",
            output=images,
            metadata={"image_count": len(images), "source_format": "docx"},
        )
    except Exception as e:
        logger.debug(f"Failed to parse {path}: {e}")
        return ParseResult.failed(str(e), modality="image")


registry.register([".docx", ".doc"], "image", parse_docx_image)


def _docx_cell_text(cell) -> str:
    """Flatten a DOCX table cell to scalar plain text.

    A cell may contain multiple paragraphs (and, rarely, nested tables — whose
    text python-docx still surfaces via ``cell.text``). Join runs on spaces and
    collapse internal whitespace so DataFrame cells stay scalar.
    """
    text = cell.text if cell.text is not None else ""
    return " ".join(text.split())


def _docx_table_to_dataframe(table, limit: int):
    """Convert one DOCX table to a DataFrame, first row as header.

    Handles ragged rows (merged cells make python-docx report uneven column
    counts) by padding/truncating each row to the header width. Returns None
    for an empty table so callers can skip it.
    """
    import pandas as pd

    rows = table.rows
    if not rows:
        return None

    # De-duplicate blank/repeated header names so pandas does not choke on
    # duplicate columns (merged header cells repeat their text).
    raw_header = [_docx_cell_text(c) for c in rows[0].cells]
    header = []
    seen: dict[str, int] = {}
    for i, name in enumerate(raw_header):
        base = name or f"col_{i}"
        if base in seen:
            seen[base] += 1
            base = f"{base}_{seen[base]}"
        else:
            seen[base] = 0
        header.append(base)

    width = len(header)
    if width == 0:
        return None

    data = []
    for row in rows[1:]:
        if len(data) >= limit:
            break
        values = [_docx_cell_text(c) for c in row.cells]
        if len(values) < width:
            values = values + [""] * (width - len(values))
        elif len(values) > width:
            values = values[:width]
        data.append(values)

    df = pd.DataFrame(data, columns=header)
    if df.empty and not any(str(h).strip() for h in header):
        return None
    return df


def parse_docx_tables(path: str, config: dict, services: dict = None) -> ParseResult:
    """Extract every table from a .docx as a dict of DataFrames.

    Complements ``parse_docx_text``, which flags ``also_contains=["tabular"]``
    when a document has tables. Output mirrors the tabular package's shape so
    the ``textualize_tabular`` task consumes it unchanged.
    """
    try:
        import pandas as pd  # noqa: F401  (used by _docx_table_to_dataframe)
        from docx import Document
    except ImportError as e:
        logger.debug(f"Missing dependency: {e}")
        return ParseResult.failed(f"Missing dependency: {e}", modality="tabular")

    try:
        doc = Document(path)
    except Exception as e:
        logger.debug(f"Failed to open {path}: {e}")
        return ParseResult.failed(str(e), modality="tabular")

    limit = (config or {}).get("max_rows", DEFAULT_MAX_TABLE_ROWS)
    all_tables = {}
    for idx, table in enumerate(doc.tables):
        try:
            df = _docx_table_to_dataframe(table, limit)
        except Exception as e:
            logger.debug(f"Skipping table {idx} in {Path(path).name}: {e}")
            continue
        if df is None or df.empty:
            continue
        all_tables[f"table_{idx + 1}"] = df

    if not all_tables:
        return ParseResult.failed(
            "No extractable tables found in DOCX", modality="tabular"
        )

    table_meta = {}
    total_rows = 0
    for name, df in all_tables.items():
        total_rows += len(df)
        table_meta[name] = {
            "row_count": len(df),
            "column_count": len(df.columns),
            "columns": list(df.columns),
            "dtypes": {col: str(dtype) for col, dtype in df.dtypes.items()},
        }

    logger.debug(
        f"DOCX tables parsed: {Path(path).name} — "
        f"{len(all_tables)} table(s), {total_rows} rows"
    )
    return ParseResult(
        modality="tabular",
        output=all_tables,
        metadata={
            "total_rows": total_rows,
            "table_count": len(all_tables),
            "table_names": list(all_tables.keys()),
            "tables": table_meta,
            "source_format": "docx",
        },
    )


registry.register([".docx", ".doc"], "tabular", parse_docx_tables)


# ===================================================================
# PPTX
# ===================================================================

def parse_pptx_text(path: str, config: dict, services: dict = None) -> ParseResult:
    """Extract text from a PowerPoint. Detects embedded images."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.debug("python-pptx not installed")
        return ParseResult.failed("python-pptx not installed", modality="text")

    try:
        t0 = time.time()
        limit = max_chars(config)
        prs = Presentation(path)

        text_runs = []
        current_len = 0
        image_count = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_runs.append(shape.text)
                    current_len += len(shape.text)
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    image_count += 1
            if current_len > limit:
                break

        content = clean_text("\n".join(text_runs)[:limit])

        also_contains = []
        if image_count > 0:
            also_contains.append("image")

        logger.debug(
            f"PPTX parsed: {Path(path).name} — {len(prs.slides)} slides, "
            f"{len(content)} chars in {time.time() - t0:.2f}s"
        )
        return ParseResult(
            modality="text",
            output=content,
            metadata={
                "char_count": len(content),
                "slide_count": len(prs.slides),
                "image_count": image_count,
                "has_images": image_count > 0,
            },
            also_contains=also_contains,
        )
    except Exception as e:
        logger.debug(f"Failed to parse {path}: {e}")
        return ParseResult.failed(str(e), modality="text")


registry.register(".pptx", "text", parse_pptx_text)


def parse_pptx_image(path: str, config: dict, services: dict = None) -> ParseResult:
    """Extract embedded images from a PPTX as PIL.Image objects."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from PIL import Image
        import io
    except ImportError as e:
        logger.debug(f"Missing dependency: {e}")
        return ParseResult.failed(f"Missing dependency: {e}", modality="image")

    try:
        prs = Presentation(path)
        images = []
        max_images = config.get("max_images", 50)

        for slide in prs.slides:
            if len(images) >= max_images:
                break
            for shape in slide.shapes:
                if len(images) >= max_images:
                    break
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_data = shape.image.blob
                        img = Image.open(io.BytesIO(image_data))
                        images.append(img)
                    except Exception:
                        continue

        if not images:
            return ParseResult.failed(
                "No extractable images found in PPTX", modality="image"
            )

        return ParseResult(
            modality="image",
            output=images,
            metadata={"image_count": len(images), "source_format": "pptx"},
        )
    except Exception as e:
        logger.debug(f"Failed to parse {path}: {e}")
        return ParseResult.failed(str(e), modality="image")


registry.register(".pptx", "image", parse_pptx_image)
